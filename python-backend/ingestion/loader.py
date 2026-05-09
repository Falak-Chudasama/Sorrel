import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from pathlib import Path

def load_document(file_path: str) -> list[dict]:
    path = Path(file_path)
    ext = path.suffix.lower()
    if ext == ".pdf":
        return _load_pdf(file_path)
    elif ext == ".docx":
        return _load_docx(file_path)
    elif ext == ".pptx":
        return _load_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

def _load_pdf(file_path: str) -> list[dict]:
    sections = []
    source_file = Path(file_path).name
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if not text or not text.strip():
                continue
            lines = [l.strip() for l in text.split("\n") if l.strip()]
            heading = lines[0] if lines else f"Page {i + 1}"
            sections.append({
                "text": text.strip(),
                "source_file": source_file,
                "doc_type": "pdf",
                "section_heading": heading[:100],
                "page_number": i + 1
            })
    return sections

def _load_docx(file_path: str) -> list[dict]:
    doc = DocxDocument(file_path)
    source_file = Path(file_path).name
    sections = []
    current_heading = "Introduction"
    current_text_parts = []
    page_number = 1
    
    for para in doc.paragraphs:
        style_name = para.style.name.lower()
        text = para.text.strip()
        if not text:
            continue
        if "heading" in style_name:
            if current_text_parts:
                sections.append({
                    "text": "\n".join(current_text_parts),
                    "source_file": source_file,
                    "doc_type": "docx",
                    "section_heading": current_heading,
                    "page_number": page_number
                })
            current_heading = text[:100]
            current_text_parts = []
            page_number += 1
        else:
            current_text_parts.append(text)
            
    if current_text_parts:
        sections.append({
            "text": "\n".join(current_text_parts),
            "source_file": source_file,
            "doc_type": "docx",
            "section_heading": current_heading,
            "page_number": page_number
        })
    return sections

def _load_pptx(file_path: str) -> list[dict]:
    prs = Presentation(file_path)
    source_file = Path(file_path).name
    sections = []
    for i, slide in enumerate(prs.slides):
        slide_title = ""
        slide_body_parts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if shape.shape_type == 13 or (hasattr(shape, "placeholder_format") and shape.placeholder_format is not None and shape.placeholder_format.idx == 0):
                slide_title = text[:100]
            else:
                slide_body_parts.append(text)
        full_text = "\n".join(slide_body_parts)
        if not full_text.strip() and not slide_title:
            continue
        sections.append({
            "text": (slide_title + "\n" + full_text).strip(),
            "source_file": source_file,
            "doc_type": "pptx",
            "section_heading": slide_title or f"Slide {i + 1}",
            "page_number": i + 1
        })
    return sections